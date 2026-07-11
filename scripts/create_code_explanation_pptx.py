from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "code_explanation_presentation.pptx"


BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(221, 235, 247)
GRAY = RGBColor(89, 89, 89)
DARK = RGBColor(34, 34, 34)
WHITE = RGBColor(255, 255, 255)


def set_text(run, size: int = 18, bold: bool = False, color: RGBColor = DARK) -> None:
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_text(p.runs[0], size=24, bold=True, color=BLUE)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.82), Inches(12.0), Inches(0.32))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        set_text(p2.runs[0], size=11, color=GRAY)


def add_footer(slide) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.05), Inches(12.25), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(160, 160, 160)
    line.line.color.rgb = RGBColor(160, 160, 160)
    foot = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(12.25), Inches(0.25))
    p = foot.text_frame.paragraphs[0]
    p.text = "AI-final-report code explanation | https://github.com/IssaIssa-tech/AI-final-report"
    set_text(p.runs[0], size=8, color=GRAY)
    p.alignment = PP_ALIGN.CENTER


def add_bullets(slide, items: list[str], x=0.8, y=1.35, w=6.0, h=5.2, font_size=16) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        set_text(p.runs[0], size=font_size)


def add_panel(slide, title: str, body: list[str], x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = BLUE
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    set_text(p.runs[0], size=15, bold=True, color=BLUE)
    for item in body:
        bp = tf.add_paragraph()
        bp.text = item
        bp.level = 0
        bp.space_after = Pt(3)
        set_text(bp.runs[0], size=10)


def add_code_box(slide, code: str, x=7.0, y=1.35, w=5.6, h=4.8) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    box.line.color.rgb = RGBColor(180, 180, 180)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = code
    for run in p.runs:
        run.font.name = "Consolas"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(40, 40, 40)


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.7), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "Source Code Explanation"
    set_text(p.runs[0], size=38, bold=True, color=BLUE)
    subtitle = slide.shapes.add_textbox(Inches(0.85), Inches(2.75), Inches(11.2), Inches(1.0))
    p2 = subtitle.text_frame.paragraphs[0]
    p2.text = "AI-final-report: ChinaXiv-style survey paper reproducibility package"
    set_text(p2.runs[0], size=20, color=DARK)
    info = slide.shapes.add_textbox(Inches(0.85), Inches(4.0), Inches(11.2), Inches(1.1))
    p3 = info.text_frame.paragraphs[0]
    p3.text = "Presenter: ISSA ISSA RASHID | Repository: https://github.com/IssaIssa-tech/AI-final-report"
    set_text(p3.runs[0], size=15, color=GRAY)
    add_footer(slide)


def slide_overview(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "1. What the Code Does", "The code supports reproducibility for a survey paper, not model training.")
    add_bullets(
        slide,
        [
            "Validates that all citations in paper.md exist in references.bib.",
            "Generates Markdown taxonomy tables from CSV source data.",
            "Generates academic grayscale figures from the taxonomy CSV files.",
            "Formats the manuscript into ChinaXiv-style Markdown and Word files.",
            "Compares the final PDF layout against the provided ChinaXiv sample.",
        ],
        w=5.8,
    )
    add_code_box(
        slide,
        "python scripts/validate_references.py\n"
        "python scripts/generate_tables.py\n"
        "python scripts/generate_figures.py\n"
        "python scripts/format_chinaxiv.py\n"
        "python scripts/compare_chinaxiv_layout.py",
    )
    add_footer(slide)


def slide_workflow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "2. Reproducible Workflow", "Each output is generated from editable source files.")
    panels = [
        ("Inputs", ["paper.md", "references.bib", "model_taxonomy.csv", "benchmark_taxonomy.csv"]),
        ("Validation", ["validate_references.py", "checks citation keys", "prevents missing references"]),
        ("Generation", ["generate_tables.py", "generate_figures.py", "format_chinaxiv.py"]),
        ("Outputs", ["paper_chinaxiv.md", "paper_chinaxiv.docx", "final PDF", "comparison report"]),
    ]
    x = 0.55
    for title, body in panels:
        add_panel(slide, title, body, x, 1.65, 2.9, 2.4)
        x += 3.05
    add_bullets(
        slide,
        [
            "Main defense point: the paper is reproducible because its tables, figures, citations, and formatting can be regenerated.",
            "The scripts do not invent experimental scores; they transform and verify documented survey materials.",
        ],
        x=0.8,
        y=4.65,
        w=11.4,
        h=1.0,
        font_size=15,
    )
    add_footer(slide)


def slide_validate_refs(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "3. validate_references.py", "Purpose: make sure every cited paper has a bibliography entry.")
    add_bullets(
        slide,
        [
            "Reads paper/paper.md and paper/references.bib.",
            "Uses regular expressions to find citation keys like [@radford2021clip].",
            "Uses another regular expression to extract BibTeX keys from references.bib.",
            "Fails with exit code 1 if a cited key is missing.",
            "Current check: 42 citation keys and 42 bibliography entries.",
        ],
        w=6.0,
    )
    add_code_box(
        slide,
        "CITATION_PATTERN = re.compile(r\"\\[@([A-Za-z0-9:_-]+)\\]\")\n"
        "BIBKEY_PATTERN = re.compile(r\"@\\w+\\{([^,]+),\")\n\n"
        "missing = sorted(cited_keys - bib_keys)\n"
        "if missing:\n"
        "    raise SystemExit(1)",
    )
    add_footer(slide)


def slide_tables(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "4. generate_tables.py", "Purpose: regenerate survey tables from CSV files.")
    add_bullets(
        slide,
        [
            "Inputs are structured CSV files in paper/tables/.",
            "The script uses Python csv.DictReader instead of manual string parsing.",
            "Outputs Markdown tables to results/tables/.",
            "This makes table updates auditable: edit CSV, rerun script, compare output.",
        ],
        w=6.0,
    )
    add_code_box(
        slide,
        "for csv_path in sorted(PAPER_TABLES.glob(\"*.csv\")):\n"
        "    rows = read_csv(csv_path)\n"
        "    output_path = RESULT_TABLES / f\"{csv_path.stem}.md\"\n"
        "    write_markdown_table(rows, output_path)",
    )
    add_footer(slide)


def slide_figures(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "5. generate_figures.py", "Purpose: draw standard academic charts from survey metadata.")
    add_bullets(
        slide,
        [
            "Reads the model and benchmark CSV files.",
            "Counts representative items by publication year using Counter.",
            "Creates grayscale, print-friendly charts with Matplotlib.",
            "Saves 300 DPI images into paper/figures/.",
            "The charts are based on transparent metadata, not fabricated experiment results.",
        ],
        w=6.0,
    )
    add_code_box(
        slide,
        "counts = Counter(row[\"year\"] for row in rows)\n"
        "years = sorted(counts)\n"
        "values = [counts[year] for year in years]\n\n"
        "plt.savefig(FIG_DIR / \"model_timeline.png\", dpi=300)",
    )
    add_footer(slide)


def slide_format(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "6. format_chinaxiv.py", "Purpose: produce the final ChinaXiv-style manuscript files.")
    add_bullets(
        slide,
        [
            "Converts draft citation keys into numbered citations.",
            "Formats references in a GB/T 7714-like numbered style.",
            "Inserts compact taxonomy tables and generated figures.",
            "Creates paper_chinaxiv.md and paper_chinaxiv.docx.",
            "Sets Letter page size, narrow margins, 10 pt body text, 12 pt headings, and a footer rule.",
        ],
        w=6.0,
    )
    add_code_box(
        slide,
        "order = citation_order(original)\n"
        "text = replace_citations(original, order)\n\n"
        "section.page_width = Cm(21.59)\n"
        "section.page_height = Cm(27.94)\n"
        "set_paragraph_top_border(footer)",
    )
    add_footer(slide)


def slide_layout(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "7. compare_chinaxiv_layout.py", "Purpose: compare final PDF layout with the ChinaXiv sample.")
    add_bullets(
        slide,
        [
            "Uses PyMuPDF to read PDF page dimensions, text boxes, font sizes, and fonts.",
            "Compares the provided sample against the final PDF.",
            "Writes docs/chinaxiv_sample_comparison.md.",
            "Verified page size: 612 x 792 pt, matching the sample.",
            "Verified side margins are close to the sample: about 133 pt.",
        ],
        w=6.0,
    )
    add_code_box(
        slide,
        "doc = fitz.open(path)\n"
        "page_width_pt = round(page.rect.width, 1)\n"
        "left = min(span[\"bbox\"][0] for span in spans)\n"
        "right = width - max(span[\"bbox\"][2] for span in spans)",
    )
    add_footer(slide)


def slide_no_training(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "8. Why There Is No Training Code", "This is an important defense answer.")
    add_bullets(
        slide,
        [
            "The project is a survey paper, not an original model-training experiment.",
            "Adding fake training or fake benchmark numbers would violate the paper instructions.",
            "The included code supports reproducibility of the survey artifacts.",
            "If original experiments are added later, src/training, src/evaluation, and configs must be expanded with real model code, seeds, logs, and hardware details.",
        ],
        w=11.5,
        font_size=17,
    )
    add_footer(slide)


def slide_repository(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "9. Repository Structure", "Where to find the code and outputs.")
    add_panel(slide, "scripts/", ["all executable reproducibility scripts", "formatting, figures, tables, validation"], 0.7, 1.35, 3.8, 1.7)
    add_panel(slide, "paper/", ["paper.md source", "references.bib", "final DOCX/PDF", "tables and figures"], 4.8, 1.35, 3.8, 1.7)
    add_panel(slide, "src/", ["reserved package structure", "models, preprocessing, training, evaluation, inference"], 8.9, 1.35, 3.8, 1.7)
    add_panel(slide, "docs/", ["layout comparison report", "this PowerPoint presentation"], 0.7, 3.45, 3.8, 1.55)
    add_panel(slide, "configs/", ["literature_review.yaml", "paper metadata and policy flags"], 4.8, 3.45, 3.8, 1.55)
    add_panel(slide, "environment", ["requirements.txt", "environment.yml", "reproducible dependencies"], 8.9, 3.45, 3.8, 1.55)
    add_footer(slide)


def slide_questions_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "10. Likely Questions and Answers", "Use these short answers in your defense.")
    questions = [
        "Q: What programming language did you use?\nA: Python. The scripts use python-docx, matplotlib, PyMuPDF, csv, pathlib, and regular expressions.",
        "Q: What is the most important script?\nA: format_chinaxiv.py, because it converts the draft into the final ChinaXiv-style manuscript.",
        "Q: How did you avoid fake references?\nA: validate_references.py checks that every citation key in paper.md exists in references.bib.",
        "Q: Are the figures real?\nA: Yes. They are generated from the model and benchmark taxonomy CSV files, not invented results.",
    ]
    add_bullets(slide, questions, x=0.7, y=1.35, w=12.0, h=5.4, font_size=13)
    add_footer(slide)


def slide_questions_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "11. More Defense Questions", "Answers for reproducibility and limitations.")
    questions = [
        "Q: Why did you not train an MLLM?\nA: The assignment output is a survey paper. Training would require datasets, model checkpoints, compute logs, seeds, and executed results.",
        "Q: How can another person reproduce your outputs?\nA: Install requirements, then run validate_references.py, generate_tables.py, generate_figures.py, format_chinaxiv.py, and compare_chinaxiv_layout.py.",
        "Q: What did the layout comparison check?\nA: Page size, approximate margins, font sizes, and common PDF fonts against the provided ChinaXiv sample.",
        "Q: What is the GitHub link?\nA: https://github.com/IssaIssa-tech/AI-final-report",
    ]
    add_bullets(slide, questions, x=0.7, y=1.35, w=12.0, h=5.4, font_size=13)
    add_footer(slide)


def slide_commands(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "12. Commands to Demonstrate", "These commands prove the source code is usable.")
    add_code_box(
        slide,
        "pip install -r requirements.txt\n\n"
        "python scripts/validate_references.py\n"
        "python scripts/generate_tables.py\n"
        "python scripts/generate_figures.py\n"
        "python scripts/format_chinaxiv.py\n"
        "python scripts/compare_chinaxiv_layout.py\n\n"
        "git log --oneline -3",
        x=0.85,
        y=1.35,
        w=11.7,
        h=4.6,
    )
    add_bullets(slide, ["Expected validation result: 42 citation keys, 42 bibliography entries, reference validation passed."], x=0.85, y=6.0, w=11.7, h=0.55, font_size=14)
    add_footer(slide)


def slide_close(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "13. Final Defense Message", "One clear statement to finish with.")
    add_bullets(
        slide,
        [
            "This codebase does not claim artificial training results.",
            "It makes the survey reproducible by generating and checking the bibliography, tables, figures, formatting, and layout comparison.",
            "The final paper, source code, figures, and validation scripts are available in the GitHub repository.",
        ],
        x=1.0,
        y=1.8,
        w=11.0,
        h=3.0,
        font_size=20,
    )
    add_footer(slide)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_overview(prs)
    slide_workflow(prs)
    slide_validate_refs(prs)
    slide_tables(prs)
    slide_figures(prs)
    slide_format(prs)
    slide_layout(prs)
    slide_no_training(prs)
    slide_repository(prs)
    slide_questions_1(prs)
    slide_questions_2(prs)
    slide_commands(prs)
    slide_close(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()

